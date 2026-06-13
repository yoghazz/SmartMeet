import io
import logging
from datetime import datetime
from typing import Dict, Any, List, Optional

logger = logging.getLogger(__name__)

# Color palette from PRD
PRIMARY_BLUE = (26, 86, 166)       # #1A56A6
SECONDARY_BLUE = (46, 134, 171)    # #2E86AB
ACCENT_ORANGE = (244, 162, 97)     # #F4A261
WHITE = (255, 255, 255)
DARK_GRAY = (50, 50, 50)
LIGHT_GRAY = (240, 240, 240)


class DocumentService:

    # ─── PDF ───────────────────────────────────────────────────────────────────

    def generate_pdf(
        self,
        title: str,
        summary: str,
        key_points: List[str],
        action_items: List[Dict],
        conclusions: str,
        transcript: str,
        participants: Optional[List[str]] = None,
        keywords: Optional[List[str]] = None,
        theme: str = "professional",
    ) -> bytes:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import cm
        from reportlab.lib import colors
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle,
            HRFlowable, PageBreak
        )
        from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY

        buffer = io.BytesIO()
        doc = SimpleDocTemplate(
            buffer,
            pagesize=A4,
            leftMargin=2.5 * cm,
            rightMargin=2.5 * cm,
            topMargin=2.5 * cm,
            bottomMargin=2.5 * cm,
        )

        rl_primary = colors.Color(PRIMARY_BLUE[0]/255, PRIMARY_BLUE[1]/255, PRIMARY_BLUE[2]/255)
        rl_accent = colors.Color(ACCENT_ORANGE[0]/255, ACCENT_ORANGE[1]/255, ACCENT_ORANGE[2]/255)
        rl_light = colors.Color(0.94, 0.94, 0.94)

        styles = getSampleStyleSheet()
        style_title = ParagraphStyle("cover_title", fontSize=28, textColor=rl_primary,
                                     alignment=TA_CENTER, spaceAfter=12, fontName="Helvetica-Bold")
        style_h1 = ParagraphStyle("h1", fontSize=16, textColor=rl_primary,
                                  spaceBefore=18, spaceAfter=8, fontName="Helvetica-Bold")
        style_h2 = ParagraphStyle("h2", fontSize=13, textColor=colors.Color(0.18, 0.33, 0.65),
                                  spaceBefore=12, spaceAfter=6, fontName="Helvetica-Bold")
        style_body = ParagraphStyle("body", fontSize=11, leading=16, spaceAfter=6,
                                    alignment=TA_JUSTIFY, fontName="Helvetica")
        style_bullet = ParagraphStyle("bullet", fontSize=11, leading=16, leftIndent=20,
                                      bulletIndent=10, spaceAfter=4, fontName="Helvetica")
        style_meta = ParagraphStyle("meta", fontSize=10, textColor=colors.gray,
                                    alignment=TA_CENTER, spaceAfter=4)

        story = []

        # ── Cover ──
        story.append(Spacer(1, 3 * cm))
        story.append(Paragraph("SmartMeet", style_meta))
        story.append(Paragraph(title, style_title))
        story.append(Spacer(1, 0.5 * cm))
        story.append(HRFlowable(width="80%", thickness=3, color=rl_accent, hAlign="CENTER"))
        story.append(Spacer(1, 0.5 * cm))

        date_str = datetime.now().strftime("%d %B %Y")
        story.append(Paragraph(f"Tanggal: {date_str}", style_meta))
        if participants:
            story.append(Paragraph(f"Peserta: {', '.join(participants)}", style_meta))
        story.append(PageBreak())

        # ── Rangkuman Eksekutif ──
        story.append(Paragraph("Rangkuman Eksekutif", style_h1))
        story.append(HRFlowable(width="100%", thickness=1, color=rl_primary))
        story.append(Spacer(1, 0.3 * cm))
        for para in summary.split("\n\n"):
            if para.strip():
                story.append(Paragraph(para.strip(), style_body))
        story.append(Spacer(1, 0.5 * cm))

        # ── Poin-Poin Kunci ──
        if key_points:
            story.append(Paragraph("Poin-Poin Kunci", style_h1))
            story.append(HRFlowable(width="100%", thickness=1, color=rl_primary))
            story.append(Spacer(1, 0.3 * cm))
            for point in key_points:
                story.append(Paragraph(f"• {point}", style_bullet))
            story.append(Spacer(1, 0.5 * cm))

        # ── Action Items ──
        if action_items:
            story.append(Paragraph("Action Items", style_h1))
            story.append(HRFlowable(width="100%", thickness=1, color=rl_primary))
            story.append(Spacer(1, 0.3 * cm))

            table_data = [["#", "Tugas", "PIC", "Tenggat"]]
            for i, item in enumerate(action_items, 1):
                table_data.append([
                    str(i),
                    item.get("task", ""),
                    item.get("assignee") or "-",
                    item.get("due_date") or "-",
                ])

            tbl = Table(table_data, colWidths=[1 * cm, 9 * cm, 4 * cm, 3.5 * cm])
            tbl.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), rl_primary),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("FONTSIZE", (0, 0), (-1, -1), 10),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, rl_light]),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.lightgrey),
                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ("TOPPADDING", (0, 0), (-1, -1), 6),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
            ]))
            story.append(tbl)
            story.append(Spacer(1, 0.5 * cm))

        # ── Kesimpulan ──
        if conclusions:
            story.append(Paragraph("Kesimpulan & Rekomendasi", style_h1))
            story.append(HRFlowable(width="100%", thickness=1, color=rl_primary))
            story.append(Spacer(1, 0.3 * cm))
            story.append(Paragraph(conclusions, style_body))
            story.append(Spacer(1, 0.5 * cm))

        # ── Transkripsi ──
        if transcript:
            story.append(PageBreak())
            story.append(Paragraph("Transkripsi Lengkap", style_h1))
            story.append(HRFlowable(width="100%", thickness=1, color=rl_primary))
            story.append(Spacer(1, 0.3 * cm))
            style_transcript = ParagraphStyle("transcript", fontSize=9, leading=14,
                                              spaceAfter=4, fontName="Helvetica")
            for line in transcript.split("\n"):
                if line.strip():
                    story.append(Paragraph(line.strip(), style_transcript))

        doc.build(story)
        return buffer.getvalue()

    # ─── DOCX ──────────────────────────────────────────────────────────────────

    def generate_docx(
        self,
        title: str,
        summary: str,
        key_points: List[str],
        action_items: List[Dict],
        conclusions: str,
        transcript: str,
        participants: Optional[List[str]] = None,
        keywords: Optional[List[str]] = None,
        theme: str = "professional",
    ) -> bytes:
        from docx import Document as DocxDocument
        from docx.shared import Pt, RGBColor, Inches, Cm
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement

        doc = DocxDocument()

        # Page margins
        for section in doc.sections:
            section.top_margin = Cm(2.5)
            section.bottom_margin = Cm(2.5)
            section.left_margin = Cm(3)
            section.right_margin = Cm(3)

        def set_heading_color(paragraph, color_rgb):
            for run in paragraph.runs:
                run.font.color.rgb = RGBColor(*color_rgb)

        def add_heading(text: str, level: int, color=PRIMARY_BLUE):
            h = doc.add_heading(text, level=level)
            set_heading_color(h, color)
            return h

        def add_horizontal_rule():
            p = doc.add_paragraph()
            pPr = p._p.get_or_add_pPr()
            pBdr = OxmlElement("w:pBdr")
            bottom = OxmlElement("w:bottom")
            bottom.set(qn("w:val"), "single")
            bottom.set(qn("w:sz"), "6")
            bottom.set(qn("w:space"), "1")
            bottom.set(qn("w:color"), f"{PRIMARY_BLUE[0]:02X}{PRIMARY_BLUE[1]:02X}{PRIMARY_BLUE[2]:02X}")
            pBdr.append(bottom)
            pPr.append(pBdr)

        # ── Cover ──
        cover_title = doc.add_heading(title, level=0)
        cover_title.alignment = WD_ALIGN_PARAGRAPH.CENTER
        set_heading_color(cover_title, PRIMARY_BLUE)

        meta = doc.add_paragraph(f"SmartMeet | {datetime.now().strftime('%d %B %Y')}")
        meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
        meta.runs[0].font.color.rgb = RGBColor(128, 128, 128)

        if participants:
            p = doc.add_paragraph(f"Peserta: {', '.join(participants)}")
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_page_break()

        # ── Rangkuman Eksekutif ──
        add_heading("Rangkuman Eksekutif", 1)
        add_horizontal_rule()
        for para in summary.split("\n\n"):
            if para.strip():
                p = doc.add_paragraph(para.strip())
                p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

        # ── Poin Kunci ──
        if key_points:
            add_heading("Poin-Poin Kunci", 1)
            add_horizontal_rule()
            for point in key_points:
                doc.add_paragraph(point, style="List Bullet")

        # ── Action Items ──
        if action_items:
            add_heading("Action Items", 1)
            add_horizontal_rule()
            table = doc.add_table(rows=1, cols=4)
            table.style = "Light Shading Accent 1"
            headers = ["#", "Tugas", "PIC", "Tenggat"]
            for i, h in enumerate(headers):
                cell = table.rows[0].cells[i]
                cell.text = h
                cell.paragraphs[0].runs[0].font.bold = True

            for idx, item in enumerate(action_items, 1):
                row = table.add_row().cells
                row[0].text = str(idx)
                row[1].text = item.get("task", "")
                row[2].text = item.get("assignee") or "-"
                row[3].text = item.get("due_date") or "-"

        # ── Kesimpulan ──
        if conclusions:
            add_heading("Kesimpulan & Rekomendasi", 1)
            add_horizontal_rule()
            p = doc.add_paragraph(conclusions)
            p.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY

        # ── Transkripsi ──
        if transcript:
            doc.add_page_break()
            add_heading("Transkripsi Lengkap", 1)
            add_horizontal_rule()
            for line in transcript.split("\n"):
                if line.strip():
                    p = doc.add_paragraph(line.strip())
                    p.runs[0].font.size = Pt(9) if p.runs else None

        buffer = io.BytesIO()
        doc.save(buffer)
        return buffer.getvalue()

    # ─── PPTX ──────────────────────────────────────────────────────────────────

    def generate_pptx(
        self,
        title: str,
        summary: str,
        key_points: List[str],
        action_items: List[Dict],
        conclusions: str,
        slides_content: Optional[List[Dict]] = None,
        participants: Optional[List[str]] = None,
        theme: str = "professional",
    ) -> bytes:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor as PPTXColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        blank_layout = prs.slide_layouts[6]  # blank

        def rgb(r, g, b):
            return PPTXColor(r, g, b)

        def add_slide(slide_title: str, bullets: List[str], is_cover: bool = False):
            slide = prs.slides.add_slide(blank_layout)
            width = prs.slide_width
            height = prs.slide_height

            # Header bar
            from pptx.util import Inches
            bg = slide.shapes.add_shape(1, 0, 0, width, Inches(1.4))  # MSO_SHAPE.RECTANGLE = 1
            bg.fill.solid()
            bg.fill.fore_color.rgb = rgb(*PRIMARY_BLUE)
            bg.line.fill.background()

            # Title text
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), width - Inches(1), Inches(1))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_title
            p.font.size = Pt(28 if is_cover else 22)
            p.font.bold = True
            p.font.color.rgb = rgb(*WHITE)

            if is_cover:
                # Subtitle
                sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), width - Inches(1), Inches(1))
                sub_tf = sub_box.text_frame
                sub_p = sub_tf.paragraphs[0]
                sub_p.text = f"{datetime.now().strftime('%d %B %Y')}"
                if participants:
                    sub_p.text += f" | {', '.join(participants[:3])}"
                sub_p.font.size = Pt(14)
                sub_p.font.color.rgb = rgb(*SECONDARY_BLUE)
                return slide

            # Content area
            content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), width - Inches(1.2), height - Inches(2.2))
            ctf = content_box.text_frame
            ctf.word_wrap = True

            for i, bullet in enumerate(bullets):
                bp = ctf.paragraphs[0] if i == 0 else ctf.add_paragraph()
                bp.text = bullet
                bp.font.size = Pt(16)
                bp.font.color.rgb = rgb(*DARK_GRAY)
                bp.level = 0

            # Accent line at bottom
            accent = slide.shapes.add_shape(1, 0, height - Inches(0.15), width, Inches(0.15))
            accent.fill.solid()
            accent.fill.fore_color.rgb = rgb(*ACCENT_ORANGE)
            accent.line.fill.background()

            return slide

        # ── Cover slide ──
        add_slide(title, [], is_cover=True)

        # ── Summary slide ──
        summary_bullets = [s.strip() for s in summary.split(". ") if s.strip()][:5]
        add_slide("Ringkasan", summary_bullets)

        # ── Key points slides ──
        if key_points:
            chunk_size = 5
            for i in range(0, len(key_points), chunk_size):
                chunk = key_points[i:i + chunk_size]
                add_slide("Poin-Poin Kunci", [f"• {p}" for p in chunk])

        # ── Custom AI slides if available ──
        if slides_content:
            for slide_data in slides_content:
                add_slide(
                    slide_data.get("title", ""),
                    slide_data.get("bullets", []),
                )

        # ── Action items slide ──
        if action_items:
            bullets = [
                f"• {item.get('task', '')} — {item.get('assignee') or 'TBD'}"
                for item in action_items[:8]
            ]
            add_slide("Action Items", bullets)

        # ── Conclusion slide ──
        if conclusions:
            conc_bullets = [s.strip() for s in conclusions.split(". ") if s.strip()][:5]
            add_slide("Kesimpulan & Langkah Selanjutnya", conc_bullets)

        buffer = io.BytesIO()
        prs.save(buffer)
        return buffer.getvalue()


document_service = DocumentService()
